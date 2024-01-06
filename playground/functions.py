from pathlib import Path

from openpyxl.cell import Cell
from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
from openpyxl.styles import Font, Border, PatternFill
from openpyxl.styles import Side
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide


# PPT에서 텍스트 추출하기
def extract_text_in_ppt(shapes, lines=None) -> list[str]:
    if lines is None:
        lines = []
    for shape in shapes:

        # 1. 일반적인 경우
        if shape.has_text_frame:
            add_line_from_text_frame(lines, shape.text_frame)

        # 2. 표인 경우
        elif shape.has_table:
            row_generator = (row for row in shape.table.rows)
            for row in row_generator:
                cell_generator = (cell for cell in row.cells)
                for cell in cell_generator:
                    add_line_from_text_frame(lines, cell.text_frame)

        # 3. 그룹인 경우 -> 재귀
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_text_in_ppt(shape.shapes, lines)
    return lines


def add_line_from_text_frame(lines: list[str], text_frame):
    # openpyxl에서 일부 문자열을 인식하지 못하는 경우가 있어 전처리
    text = ILLEGAL_CHARACTERS_RE.sub(r'', text_frame.text)
    # 긴 텍스트의 경우 100자로 잘라서 보내기
    while True:
        text = text[:100]
        if text is not None and len(text) != 0:
            lines.append(text)
        # print(text)
        text = text[100:]
        if len(text) < 100:
            break


def auto_fit_column_size(worksheet, columns=None, margin=2):
    for i, column_cells in enumerate(worksheet.columns):
        is_ok = False
        if columns is None or (isinstance(columns, list) and i in columns):
            is_ok = True

        if is_ok:
            length = max(len(str(cell.value)) for cell in column_cells)
            worksheet.column_dimensions[column_cells[0].column_letter].width = length + margin

    return worksheet


def base_border(border_style='thin', rgb='000000'):
    return Border(left=Side(border_style, rgb),
                  right=Side(border_style, rgb),
                  top=Side(border_style, rgb),
                  bottom=Side(border_style, rgb))


def base_font(name, rgb, size):
    return Font(name=name, color=rgb, size=size)


def base_fill(rgb: str):
    return PatternFill(start_color=rgb, end_color=rgb, fill_type='solid')


def design_excel(ws: Worksheet):
    rows: list[tuple[Cell]] = ws.rows

    # 폰트설정
    font_style: str = 'KoPub돋움체 Medium'

    for i, tup in enumerate(rows):
        for el in tup:
            el.fill = base_fill('000000') if i == 0 else base_fill('ffffff')
            el.font = base_font(font_style, 'ffffff', 13) if i == 0 else base_font(font_style, '000000', 11)
            el.border = base_border()
    auto_fit_column_size(ws, margin=10)


def has_digit(x):
    return any(c.isdigit() for c in x)


def read_ppt(path: Path, ws: Worksheet):
    prs: Presentation = Presentation(path)
    sls: list[Slide] = prs.slides
    for i, sl in enumerate(sls):
        print(f'[{path.name}] :: {i}번째 슬라이드 읽음')

        result_set = set()
        # for (j, shape) in enumerate(sl.shapes):
        #     result = ""
        #     if j <= 4 and shape.has_text_frame:
        #         for (k, run) in enumerate(shape.text_frame.paragraphs[0].runs):
        #             if k <= 3:
        #                 result += run.text
        #                 # print(f"[{k}] :::: {run.text}")
        #     if result != "":
        for sp in sl.shapes:
            if sp.has_text_frame:
                tf = sp.text_frame
                if sp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    text = ILLEGAL_CHARACTERS_RE.sub(r'', tf.text)
                    result_set.add(text)

        filtered: list[str] = list(filter(has_digit, result_set))
        sub = ""
        if len(filtered) > 0:
            dot_max = max(filtered, key=lambda x: len([c for c in x.split(".") if has_digit(c) and x.count(">") == 0]))
            sub = dot_max

        # PPT 텍스트 추출
        for t in extract_text_in_ppt(sl.shapes):
            append_row(i, t, path.name, ws, sub)


global gi
gi = 1


def append_row(slide_index: int, text: str, file_name: str, ws: Worksheet, sub: str):
    global gi

    # 슬라이드별 텍스트 추출 하기 & 특수문자만 있는 경우는 제외
    # regex = re.compile("[0-9a-zA-Zㄱ-힗]", re.MULTILINE)
    # if re.match(regex, text):

    # 엑셀에서 일부 특수문자는 인식이 안되는 경우가 있어 텍스트앞에 작은따옴표(') 붙이기
    text = "'" + text
    ws.append([gi, file_name, str(slide_index) + "번째", sub, text])
    gi += 1
