import re
from pathlib import Path

from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide

from playground.utils import has_digit


def read_ppt(path: Path, ws: Worksheet = None):
    prs: Presentation = Presentation(path)
    sls: list[Slide] = prs.slides
    for i, sl in enumerate(sls):
        print(f'[{path.name}] :: {i}번째 슬라이드 읽음')
        result_set = set()
        for sp in sl.shapes:
            if sp.has_text_frame:
                tf = sp.text_frame
                if sp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    text = ILLEGAL_CHARACTERS_RE.sub(r'', tf.text)
                    result_set.add(text)

        filtered: list[str] = list(filter(has_digit, result_set))
        sub = ""
        if len(filtered) > 0:
            dot_max = max(filtered, key=lambda x: len([c for c in x.split(".") if has_digit(c) and x.count(">") == 0]))
            sub = dot_max

        # PPT 텍스트 추출
        for t in __extract_text_in_ppt(sl.shapes):
            __append_row(i, t, path.name, ws, sub)


# PPT에서 텍스트 추출하기
def __extract_text_in_ppt(shapes, lines=None) -> list[str]:
    if lines is None:
        lines = []
    for shape in shapes:

        # 1. 일반적인 경우
        if shape.has_text_frame:
            __add_line_from_text_frame(lines, shape.text_frame)

        # 2. 표인 경우
        elif shape.has_table:
            row_generator = (row for row in shape.table.rows)
            for row in row_generator:
                cell_generator = (cell for cell in row.cells)
                for cell in cell_generator:
                    __add_line_from_text_frame(lines, cell.text_frame)

        # 3. 그룹인 경우 -> 재귀
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            __extract_text_in_ppt(shape.shapes, lines)
    return lines


def __add_line_from_text_frame(lines: list[str], text_frame):
    # openpyxl에서 일부 문자열을 인식하지 못하는 경우가 있어 전처리
    text = ILLEGAL_CHARACTERS_RE.sub(r'', text_frame.text)
    lines.append(text)


# 인덱스 넘버링을 위한 글로벌 변수 선언
global gi
gi = 1


def __append_row(slide_index: int, text: str, file_name: str, ws: Worksheet, sub: str):
    global gi

    # 슬라이드별 텍스트 추출 하기 & 특수문자만 있는 경우는 제외
    # regex = re.compile("[0-9a-zA-Zㄱ-힗]", re.MULTILINE)
    if text is not None and len(text) != 0 and not (re.sub(r'\W', '', text) == ''):
        # 엑셀에서 일부 특수문자는 인식이 안되는 경우가 있어 텍스트앞에 구분자(') 붙이기
        # (엑셀에서 '텍스트 나누기'로 후처리 필요함)
        text = "^" + text
        ws.append([gi, file_name, str(slide_index) + "번째", sub, text])
        gi += 1
